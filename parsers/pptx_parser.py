"""
PPTX 文件解析器
支持解析 PowerPoint 幻灯片 (.pptx)
提取文本、表格、图片占位符、备注等信息
"""

import logging
from pathlib import Path

from core.models import ExtractedElement, PageContent
from parsers import BaseParser

logger = logging.getLogger("parsers.pptx")

# 可选依赖
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt  # noqa: F401  (defensive: ensure pptx.util loads)

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx 库未安装，PPTX 解析功能不可用")

# B6/v0.11.0: 动画 XML namespace
_ANIM_NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}


class PPTXParser(BaseParser):
    """PPTX 幻灯片解析器"""

    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx", ".ppt"]

    @property
    def supported_magic(self) -> list[bytes]:
        # PPTX 是 ZIP 格式
        # PPT 是 OLE2 格式
        return [b"PK\x03\x04", b"\xd0\xcf\x11\xe0"]

    def parse(self, file_path: Path) -> list[PageContent]:
        """解析 PPTX 文件"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx 库未安装，无法解析 PPTX 文件")

        logger.info("开始解析 PPTX: %s", file_path)

        try:
            prs = Presentation(str(file_path))
        except Exception as e:
            logger.error("无法打开 PPTX 文件: %s", e)
            raise ValueError(f"无法打开 PPTX 文件: {e}") from e

        pages = []
        for idx, slide in enumerate(prs.slides, 1):
            page = self._parse_slide(slide, idx)
            pages.append(page)

        logger.info("PPTX 解析完成: %d 页", len(pages))
        return pages

    def _parse_slide(self, slide, slide_number: int) -> PageContent:
        """解析单页幻灯片"""
        elements = []
        raw_text_parts = []
        has_image = False
        has_table = False
        elem_idx = 0

        # 提取标题
        title = ""
        slide_title = getattr(slide.shapes, "title", None)
        if slide_title and hasattr(slide_title, "text") and slide_title.text.strip():
            title = slide_title.text.strip()
            elements.append(
                ExtractedElement(
                    elementId=f"elem_{slide_number}_{elem_idx}",
                    elementType="title",
                    content=title,
                    metadata={"is_title": True},
                )
            )
            raw_text_parts.append(title)
            elem_idx += 1

        # 遍历所有形状
        for shape in slide.shapes:
            # 跳过标题（已处理）
            if shape == slide_title:
                continue

            # 处理表格
            if shape.has_table:
                has_table = True
                table_text = self._extract_table(shape.table)
                if table_text:
                    elements.append(
                        ExtractedElement(
                            elementId=f"elem_{slide_number}_{elem_idx}",
                            elementType="table",
                            content=table_text,
                            metadata={"rows": len(shape.table.rows), "cols": len(shape.table.columns)},
                        )
                    )
                    raw_text_parts.append(f"[表格]\n{table_text}")
                    elem_idx += 1
                continue

            # 处理图片
            if (
                shape.shape_type is not None
                and hasattr(shape.shape_type, "name")
                and "PICTURE" in shape.shape_type.name
            ):
                has_image = True
                elements.append(
                    ExtractedElement(
                        elementId=f"elem_{slide_number}_{elem_idx}",
                        elementType="image",
                        content=f"[图片] {shape.name}",
                        metadata={
                            "shape_name": shape.name,
                            "width": shape.width.inches if hasattr(shape, "width") else None,
                            "height": shape.height.inches if hasattr(shape, "height") else None,
                        },
                    )
                )
                raw_text_parts.append(f"[图片] {shape.name}")
                elem_idx += 1
                continue

            # 处理文本
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                elem_type = self._detect_element_type(text)
                elements.append(
                    ExtractedElement(elementId=f"elem_{slide_number}_{elem_idx}", elementType=elem_type, content=text)
                )
                raw_text_parts.append(text)
                elem_idx += 1

        # 提取备注
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide:
            notes_text_frame = slide.notes_slide.notes_text_frame
            if notes_text_frame and notes_text_frame.text.strip():
                notes_text = notes_text_frame.text.strip()
                elements.append(
                    ExtractedElement(
                        elementId=f"elem_{slide_number}_{elem_idx}",
                        elementType="note",
                        content=notes_text,
                        metadata={"is_note": True},
                    )
                )
                raw_text_parts.append(f"[备注] {notes_text}")
                elem_idx += 1

        # B6/v0.11.0: 提取动画顺序（p:timing/p:par/p:animMotion/p:animEffect）
        animations = self._extract_animations(slide)

        return PageContent(
            pageNumber=slide_number,
            elements=elements,
            rawText="\n".join(raw_text_parts),
            hasImage=has_image,
            hasTable=has_table,
            metadata={
                "animations": animations,
                "animations_count": len(animations),
            },
        )

    def _extract_table(self, table) -> str:
        """提取表格文本内容"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _extract_animations(self, slide) -> list[dict]:
        """B6/v0.11.0: 提取动画顺序。

        返回 [{index, shape_id, shape_name, effect_type, delay_ms}, ...] 按播放顺序排序。
        跳过不抛异常：动画 XML 可能在某些 PPTX 里缺失/破损。
        """
        animations: list[dict] = []
        try:
            slide_elem = slide._element  # CT_Slide xml 元素
            timing = slide_elem.find(".//p:timing", _ANIM_NS)
            if timing is None:
                return animations
            time_node_list = timing.find("p:tnLst", _ANIM_NS)
            if time_node_list is None:
                return animations
            # 递归扫描 par 节点（每个 par = 一个动画触发时间线）
            idx = 0
            for par in time_node_list.iter(f"{{{_ANIM_NS['p']}}}par"):
                # delay 时长（p:stDt 内的 delay）
                delay_ms = 0
                st_cell = par.find("p:cTn/p:stCell", _ANIM_NS)
                if st_cell is not None:
                    try:
                        delay_ms = int(st_cell.get("val", "0"))
                    except (TypeError, ValueError):
                        delay_ms = 0
                # 找所有 animMotion / animEffect 子节点
                for anim_node in list(par):
                    tag = anim_node.tag.split("}")[-1] if "}" in anim_node.tag else anim_node.tag
                    if tag in ("animMotion", "animEffect", "animClr", "animScale", "animRot"):
                        # 找 shape 引用（p:spTgt/p:tgtEl）
                        shape_id = ""
                        shape_name = ""
                        sp_tgt = anim_node.find("p:cTn/p:spTgt", _ANIM_NS)
                        tgt_el = sp_tgt.find("p:tgtEl", _ANIM_NS) if sp_tgt is not None else None
                        shape_id = (tgt_el.get("spid") or "") if tgt_el is not None else ""
                        # 映射到 shape name（python-pptx 不直接给，XML 里有 nvSpPr/cNvPr name）
                        try:
                            sp = slide.shapes  # noqa
                            for s in sp:
                                s_elem = getattr(s, "_element", None)
                                if s_elem is not None and str(s_elem.get("id", "")) == shape_id:
                                    shape_name = s.name
                                    break
                        except Exception:
                            pass
                        idx += 1
                        animations.append(
                            {
                                "index": idx,
                                "shape_id": shape_id,
                                "shape_name": shape_name,
                                "effect_type": tag,
                                "delay_ms": delay_ms,
                            }
                        )
        except Exception as e:
            logger.debug("动画顺序提取失败（非致命）: %s", e)
        return animations

    def _detect_element_type(self, text: str) -> str:
        """检测元素类型"""
        text = text.strip()
        if not text:
            return "empty"

        # 检测标题
        if len(text) < 100 and (text.endswith("：") or text.endswith(":")):
            return "heading"

        # 检测列表
        if text.startswith(("•", "-", "*", "1.", "2.", "（", "(")):
            return "list"

        return "text"
