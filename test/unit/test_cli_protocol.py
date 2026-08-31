"""
CLI 协议契约测试（test_cli_protocol.py）

验证 `python -m formatforge` 的 stdout 协议 JSON 形状与退出码。
JS 侧 python-runner 依赖这些契约，改动须同步 PLUGIN_PLAN.md §4.3。
"""

import json
import os
import shutil
import subprocess
import sys
from pathlib import Path

import pytest

REPO_ROOT = Path(__file__).resolve().parents[2]
PY = sys.executable
FIXTURES = REPO_ROOT / "test" / "fixtures"


def run_cli(*args: str, stdin: str | None = None) -> tuple[dict, int]:
    env_pythonpath = str(REPO_ROOT)
    proc = subprocess.run(
        [PY, "-m", "formatforge", *args],
        capture_output=True,
        text=True,
        input=stdin,
        cwd=REPO_ROOT,
        timeout=180,
        # 未 pip install 时也能找到 formatforge 包（与 pythonpath=["."] 一致）
        env={**os.environ, "PYTHONPATH": env_pythonpath},
    )
    # stdout 首行必须是合法协议 JSON
    lines = [l for l in proc.stdout.splitlines() if l.strip()]
    assert lines, f"stdout 为空。stderr={proc.stderr[-500:]}"
    payload = json.loads(lines[0])
    assert isinstance(payload, dict)
    assert "ok" in payload and "code" in payload
    return payload, proc.returncode


class TestVersion:
    def test_version_ok(self):
        payload, code = run_cli("version")
        assert payload["ok"] is True
        assert payload["code"] == 200
        assert payload["data"]["name"] == "dsh-formatforge"
        assert code == 0


class TestFormats:
    def test_formats_lists_supported(self):
        payload, code = run_cli("formats")
        assert payload["ok"] is True
        data = payload["data"]
        assert data["count"] > 20
        for fmt in ("pdf", "docx", "xlsx", "pptx", "eml", "toml", "csv"):
            assert fmt in data["formats"]
        assert "json" in data["output_formats"]
        assert code == 0


class TestTranslateText:
    def test_stdin_text_ok(self):
        payload, code = run_cli(
            "translate", "--stdin-text", "--format", "text",
            stdin="Hello FormatForge\n第二行",
        )
        assert payload["ok"] is True
        data = payload["data"]
        assert isinstance(data["content"], str) and len(data["content"]) > 0
        assert data["format"] == "text"
        meta = data["meta"]
        assert {"parser", "file_size", "elapsed_ms"} <= set(meta)
        assert code == 0

    def test_txt_file_conversion(self):
        target = FIXTURES / "gbk_chinese.txt"
        if not target.exists():
            pytest.skip("fixture 缺失")
        payload, code = run_cli("translate", str(target), "--format", "text")
        assert payload["ok"] is True
        assert "GBK 编码测试文件" in payload["data"]["content"]
        assert code == 0


class TestTranslateErrors:
    def test_missing_file(self):
        payload, code = run_cli("translate", "/no/such/file.docx")
        assert payload["ok"] is False
        err = payload["error"]
        assert err["kind"] == "file_not_found"
        assert err["message"]
        assert code == 2

    def test_directory_rejected(self):
        payload, code = run_cli("translate", str(REPO_ROOT / "test"))
        assert payload["ok"] is False
        assert payload["error"]["kind"] == "file_not_found"
        assert code == 2


@pytest.mark.slow
class TestTranslatePdf:
    def test_pdf_with_enhance_hint(self):
        """扫描件应触发 image_only 增强提示（PLUGIN_PLAN §6）"""
        target = FIXTURES / "image_only_test.pdf"
        if not target.exists():
            pytest.skip("fixture 缺失")
        payload, code = run_cli("translate", str(target), "--format", "markdown")
        assert payload["ok"] is True
        enhance = payload["data"].get("enhance")
        assert enhance is not None
        assert enhance["needed"] is True
        assert enhance["reason"] == "image_only"
        assert "扫描件" in enhance["hint"] or "文字层" in enhance["hint"]
        assert code == 0


class TestR3SmartDefault:
    """R3.1: auto 模式下自动附带 quality（无需显式 --quality）。"""

    def test_auto_mode_implicit_quality(self):
        target = FIXTURES / "gbk_chinese.txt"
        if not target.exists():
            pytest.skip("fixture 缺失")
        # 不带 --quality；auto 模式应该自动附 quality
        payload, code = run_cli("translate", str(target), "--type", "auto", "--format", "text")
        assert payload["ok"] is True
        data = payload["data"]
        assert data.get("meta", {}).get("quality_auto") is True, "auto 模式应自动开启 quality"
        assert "quality" in data, "应附 quality 报告"
        assert "overall_score" in data["quality"]
        assert code == 0


class TestR3EncodingRetry:
    """R3.3: --encoding 透传让 quality.actions.retry_with.encoding 真可重调。"""

    def test_encoding_override_redecodes_gbk(self):
        """CLI --encoding gbk 应该用 GBK 重新解码（绕过 chardet 误判）。"""
        target = FIXTURES / "gbk_chinese.txt"
        if not target.exists():
            pytest.skip("fixture 缺失")
        # 强制 gbk 解码（fixture 本身就是 GBK）→ 中文应可读
        payload, code = run_cli(
            "translate", str(target), "--encoding", "gbk", "--format", "text",
        )
        assert payload["ok"] is True
        content = payload["data"]["content"]
        assert "GBK 编码测试文件" in content, f"GBK 重解码失败：{content[:60]!r}"
        assert code == 0

    def test_encoding_action_retry_round_trip(self):
        """模拟 actions.retry_with.encoding 闭环：first pass 不带 encoding → quality.actions;
        second pass 带 encoding → 中文可读。"""
        target = FIXTURES / "gbk_chinese.txt"
        if not target.exists():
            pytest.skip("fixture 缺失")
        # 第一轮：auto，让 chardet 自由判断
        p1, _ = run_cli("translate", str(target), "--type", "auto", "--format", "text")
        q1 = p1["data"].get("quality") or {}
        actions = q1.get("actions") or []
        encoding_actions = [a for a in actions if a.get("code") == "encoding" and a.get("retry_with")]
        # 第二轮：依 retry_with.encoding 重调
        if encoding_actions:
            enc = encoding_actions[0]["retry_with"].get("encoding")
            assert enc in ("gbk", "gb2312", "gb18030")
            p2, _ = run_cli("translate", str(target), "--encoding", enc, "--format", "text")
            assert "GBK 编码测试文件" in p2["data"]["content"]
        else:
            # chardet 在本机已经猜对了 GBK → 不出 encoding action 也属正常
            assert "GBK 编码测试文件" in p1["data"]["content"]


class TestR3MarkdownStructureField:
    """R2.3/R3.x 配套：md 格式输出应有 meta.structured 字段供会话模型识别。"""

    def test_markdown_structured_flag(self):
        target = FIXTURES / "complex_test.pdf"
        if not target.exists():
            pytest.skip("fixture 缺失")
        payload, code = run_cli("translate", str(target), "--format", "markdown")
        assert payload["ok"] is True
        meta = payload["data"]["meta"]
        # meta.structured 是 R2 引入的契约字段；测试必须守住
        assert "structured" in meta
        assert isinstance(meta["structured"], bool)
        assert code == 0


class TestR10LanguageFlag:
    """v0.10.0/B9: --language 写入 meta.target_language + enhance.hint。"""

    def test_language_sets_metadata(self):
        target = FIXTURES / "gbk_chinese.txt"
        if not target.exists():
            pytest.skip("fixture 缺失")
        payload, code = run_cli("translate", str(target), "--language", "en", "--format", "text")
        assert payload["ok"] is True
        meta = payload["data"]["meta"]
        assert meta.get("target_language") == "en"
        # enhance.hint 应提示目标语言
        enhance = payload["data"].get("enhance")
        assert enhance is not None
        assert "en" in enhance.get("hint", "")
        assert code == 0

    def test_language_case_normalized(self):
        target = FIXTURES / "gbk_chinese.txt"
        if not target.exists():
            pytest.skip("fixture 缺失")
        payload, code = run_cli("translate", str(target), "--language", "ZH-CN", "--format", "text")
        assert payload["ok"] is True
        assert payload["data"]["meta"].get("target_language") == "zh-cn"
        assert code == 0


class TestR10OutputFile:
    """v0.10.0/A9: --output-file 把 content 落盘，stdout 协议不变。"""

    def test_output_file_writes_content(self, tmp_path):
        target = FIXTURES / "gbk_chinese.txt"
        if not target.exists():
            pytest.skip("fixture 缺失")
        out_file = tmp_path / "out.txt"
        payload, code = run_cli("translate", str(target), "--format", "text", "--output-file", str(out_file))
        assert payload["ok"] is True
        # stdout 协议不变：content 仍包含转换结果
        assert "GBK 编码测试文件" in payload["data"]["content"]
        # 落盘文件存在且内容等于 content
        assert out_file.exists()
        assert out_file.read_text(encoding="utf-8") == payload["data"]["content"]
        # meta.output_file 字段标记
        assert payload["data"]["meta"].get("output_file") == str(out_file)
        assert code == 0


class TestR10FormatsCategory:
    """v0.10.0/A10: formats --category 按分类过滤。"""

    def test_category_document(self):
        payload, code = run_cli("formats", "--category", "document")
        assert payload["ok"] is True
        d = payload["data"]
        assert d["category"] == "document"
        assert d["count"] >= 5
        # 必须含主文档格式
        assert "pdf" in d["formats"]
        assert "docx" in d["formats"]
        assert "txt" in d["formats"]
        assert code == 0

    def test_category_data(self):
        payload, code = run_cli("formats", "--category", "data")
        assert payload["ok"] is True
        d = payload["data"]
        assert d["category"] == "data"
        assert "csv" in d["formats"]
        assert "xlsx" in d["formats"]
        assert code == 0

    def test_category_invalid(self):
        """argparse choices 校验在 CLI 层拒绝；非 0 退出码 + stderr 信息。"""
        proc = subprocess.run(
            [PY, "-m", "formatforge", "formats", "--category", "no_such_thing"],
            capture_output=True, text=True, cwd=REPO_ROOT, timeout=30,
            env={**os.environ, "PYTHONPATH": str(REPO_ROOT)},
        )
        assert proc.returncode != 0
        assert "invalid choice" in proc.stderr or "no_such_thing" in proc.stderr

    def test_categories_listed(self):
        payload, _ = run_cli("formats")
        cats = payload["data"].get("categories", [])
        # 必须含这 6 类
        assert {"document", "data", "email", "image", "archive", "audio"} <= set(cats)


class TestR10Batch:
    """v0.10.0/B3: batch 子命令 + --force 重转 + 报告落盘。"""

    def test_batch_basic_run(self, tmp_path):
        in_dir = tmp_path / "in"
        out_dir = tmp_path / "out"
        in_dir.mkdir()
        # 拷两个 fixture（GBK 中文 + 一个 txt）
        shutil.copy(FIXTURES / "gbk_chinese.txt", in_dir / "a.txt")
        shutil.copy(FIXTURES / "gbk_chinese.txt", in_dir / "b.txt")
        payload, code = run_cli(
            "batch", str(in_dir),
            "--out", str(out_dir),
            "--to", "markdown",
            "--workers", "2",
            "--type", "auto",
            "--force",
        )
        assert payload["ok"] is True
        assert payload["total"] == 2
        assert payload["ok_count"] == 2
        assert payload["failed"] == 0
        # 产物存在
        assert (out_dir / "a.md").exists()
        assert (out_dir / "b.md").exists()
        # 报告落盘
        report = out_dir / "_batch_report.json"
        assert report.exists()
        assert code == 0

    def test_batch_skip_existing(self, tmp_path):
        """产物比源新 → 跳过（无需 --force）。"""
        in_dir = tmp_path / "in"
        out_dir = tmp_path / "out"
        in_dir.mkdir()
        shutil.copy(FIXTURES / "gbk_chinese.txt", in_dir / "a.txt")
        # 先跑一次产产物
        run_cli("batch", str(in_dir), "--out", str(out_dir), "--to", "markdown")
        # 第二跑应全跳过
        payload, code = run_cli("batch", str(in_dir), "--out", str(out_dir), "--to", "markdown")
        assert payload["ok"] is True
        assert payload["ok_count"] == 0
        assert payload["skipped"] >= 1
        assert code == 0

    def test_batch_empty_source(self, tmp_path):
        """空目录 → 空报告（v0.10.0/B3: 但 exit != 0，与原契约一致）。"""
        empty = tmp_path / "empty"
        empty.mkdir()
        payload, code = run_cli(
            "batch", str(empty),
            "--out", str(tmp_path / "out"),
            "--to", "markdown",
            "--force",
        )
        assert payload["ok"] is True
        assert payload["total"] == 0
        assert payload["ok_count"] == 0
        # 报告必须落盘（契约）
        assert (tmp_path / "out" / "_batch_report.json").exists()
        assert code != 0  # 空源也算异常（用户期望处理但没匹配）


class TestR11CsvSchema:
    """v0.11.0/B1: CSV → structured_data 含 schema + preview_rows。"""

    def test_csv_schema_inference(self, tmp_path):
        """构造多种列类型的 CSV → schema 推断正确（integer/float/date/boolean/string）。"""
        csv_file = tmp_path / "test.csv"
        csv_file.write_text(
            "id,name,age,joined,active,salary\n"
            "1,Alice,28,2022-01-15,true,50000.5\n"
            "2,Bob,35,2021-06-30,false,75000\n"
            "3,Carol,42,2020-03-22,true,92000.75\n",
            encoding="utf-8",
        )
        payload, code = run_cli("translate", str(csv_file), "--type", "table", "--format", "json")
        assert payload["ok"] is True
        sd = payload["data"].get("structured_data") or {}
        schema = sd.get("schema", [])
        # 转 name → type 字典
        type_map = {c["name"]: c["type"] for c in schema}
        assert type_map.get("id") == "integer"
        assert type_map.get("name") == "string"
        assert type_map.get("age") == "integer"
        assert type_map.get("joined") == "date"
        assert type_map.get("active") == "boolean"
        # salary 列：50000.5/75000/92000.75 混整数+小数 → 应识别为 float（合并判定）
        assert type_map.get("salary") == "float"
        # preview_rows 应有前几行数据
        previews = sd.get("preview_rows", [])
        assert previews and len(previews[0]) >= 1
        assert code == 0

    def test_csv_pure_integers_stay_integer(self, tmp_path):
        """全整数列 → 不被识别为 float。"""
        csv_file = tmp_path / "ints.csv"
        csv_file.write_text("count\n10\n20\n30\n", encoding="utf-8")
        payload, _ = run_cli("translate", str(csv_file), "--type", "table", "--format", "json")
        sd = payload["data"].get("structured_data") or {}
        schema = sd.get("schema", [])
        assert schema and schema[0]["type"] == "integer"


class TestR11XlsxSchema:
    """v0.11.0/B1: XLSX → structured_data 含 schema。"""

    # v0.13.0: openpyxl 在某些 venv 缺失（如 py312 启动器）→ importorskip 防 ERROR
    openpyxl = pytest.importorskip("openpyxl")

    def test_xlsx_structured_data_present(self, tmp_path):
        """xlsx fixture 有表 → schema 字段非空。"""
        # 用 openpyxl 在 tmp_path 生成小 xlsx
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["id", "name", "score"])
        ws.append([1, "Alice", 95])
        ws.append([2, "Bob", 87])
        ws.append([3, "Carol", 92])
        xlsx_path = tmp_path / "mini.xlsx"
        wb.save(str(xlsx_path))
        payload, code = run_cli("translate", str(xlsx_path), "--type", "table", "--format", "json")
        assert payload["ok"] is True
        sd = payload["data"].get("structured_data") or {}
        schema = sd.get("schema", [])
        # 至少有一个 schema 条目
        assert schema and len(schema) >= 3
        # preview_rows 应有数据
        assert sd.get("preview_rows"), "preview_rows 应非空"
        assert code == 0


class TestR11DocxRevisions:
    """v0.11.0/B5: DOCX 修订追踪（w:ins/w:del）。"""

    # v0.13.0: python-docx 在某些 venv 缺失 → importorskip 防 ERROR
    docx = pytest.importorskip("docx")

    def test_docx_revisions_extracted(self, tmp_path):
        """构造带 w:ins / w:del 的 DOCX → parser 应抽到 revisions。"""
        from docx import Document
        from docx.oxml.ns import qn
        from lxml import etree

        docx_path = tmp_path / "rev.docx"
        doc = Document()
        para = doc.add_paragraph("初始段落")
        ins_xml = (
            '<w:ins xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" '
            'w:id="100" w:author="Alice" w:date="2026-08-28T10:00:00Z">'
            '<w:r><w:t>插入的内容</w:t></w:r></w:ins>'
        )
        para._p.append(etree.fromstring(ins_xml))

        para2 = doc.add_paragraph("待删除段落")
        del_xml = (
            '<w:del xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" '
            'w:id="101" w:author="Bob" w:date="2026-08-28T10:01:00Z">'
            '<w:r><w:delText>这段被删了</w:delText></w:r></w:del>'
        )
        para2._p.append(etree.fromstring(del_xml))
        doc.save(docx_path)

        # 直接调 parser（避开 CLI 协议层的 json 简化输出）
        from parsers.docx_parser import DOCXParser

        parser = DOCXParser()
        pages = parser.parse(docx_path)
        assert pages and pages[0].metadata
        revisions = pages[0].metadata.get("revisions", [])
        assert len(revisions) == 2, f"期望 2 条修订，实际 {len(revisions)}: {revisions}"

        types = {r["type"] for r in revisions}
        authors = {r["author"] for r in revisions}
        assert "ins" in types
        assert "del" in types
        assert "Alice" in authors
        assert "Bob" in authors

        # 修订文本提取
        ins_entry = next(r for r in revisions if r["type"] == "ins")
        del_entry = next(r for r in revisions if r["type"] == "del")
        assert "插入的内容" in ins_entry["text"]
        assert "这段被删了" in del_entry["text"]

    def test_docx_no_revisions_no_crash(self, tmp_path):
        """无修订的 DOCX → revisions 为空数组，不报错。"""
        from docx import Document

        docx_path = tmp_path / "plain.docx"
        doc = Document()
        doc.add_paragraph("普通段落，没有修订")
        doc.save(docx_path)

        from parsers.docx_parser import DOCXParser

        parser = DOCXParser()
        pages = parser.parse(docx_path)
        revisions = pages[0].metadata.get("revisions", [])
        assert revisions == []
        assert pages[0].metadata.get("revisions_count") == 0


class TestR11EpubChapters:
    """v0.11.0/B8: EPUB 按 toc 章节拆分，元素 metadata 带 chapter_title。"""

    @pytest.fixture
    def sample_epub(self, tmp_path):
        """生成带 NCX 的小 EPUB（2 章节）。"""
        import zipfile

        epub_path = tmp_path / "mini.epub"
        parts = {
            "mimetype": b"application/epub+zip",
            "META-INF/container.xml": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
                '<rootfiles><rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/></rootfiles>'
                "</container>"
            ).encode("utf-8"),
            "OEBPS/content.opf": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<package xmlns="http://www.idpf.org/2007/opf" version="2.0" unique-identifier="BookId">'
                '<metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
                '<dc:title>测试书</dc:title>'
                '<dc:creator>测试作者</dc:creator>'
                "<dc:identifier id=\"BookId\">urn:uuid:test</dc:identifier>"
                "</metadata>"
                '<manifest>'
                '<item id="ch1" href="chapter1.xhtml" media-type="application/xhtml+xml"/>'
                '<item id="ch2" href="chapter2.xhtml" media-type="application/xhtml+xml"/>'
                '<item id="ncx" href="toc.ncx" media-type="application/x-dtbncx+xml"/>'
                "</manifest>"
                '<spine toc="ncx"><itemref idref="ch1"/><itemref idref="ch2"/></spine>'
                "</package>"
            ).encode("utf-8"),
            "OEBPS/toc.ncx": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<ncx xmlns="http://www.daisy.org/z3986/2005/ncx/" version="2005-1">'
                '<head><meta name="dtb:uid" content="urn:uuid:test"/></head>'
                "<navMap>"
                '<navPoint id="navPoint-1" playOrder="1">'
                "<navLabel><text>第一章 引言</text></navLabel>"
                '<content src="chapter1.xhtml"/>'
                "</navPoint>"
                '<navPoint id="navPoint-2" playOrder="2">'
                "<navLabel><text>第二章 方法</text></navLabel>"
                '<content src="chapter2.xhtml"/>'
                "</navPoint>"
                "</navMap>"
                "</ncx>"
            ).encode("utf-8"),
            "OEBPS/chapter1.xhtml": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<!DOCTYPE html><html xmlns="http://www.w3.org/1999/xhtml"><body>'
                "<h1>第一章</h1><p>背景介绍。</p></body></html>"
            ).encode("utf-8"),
            "OEBPS/chapter2.xhtml": (
                '<?xml version="1.0" encoding="UTF-8"?>'
                '<!DOCTYPE html><html xmlns="http://www.w3.org/1999/xhtml"><body>'
                "<h1>第二章</h1><p>方法细节。</p></body></html>"
            ).encode("utf-8"),
        }
        with zipfile.ZipFile(epub_path, "w") as zf:
            zf.writestr(
                zipfile.ZipInfo("mimetype"), parts["mimetype"], compress_type=zipfile.ZIP_STORED
            )
            for name, content in parts.items():
                if name == "mimetype":
                    continue
                zf.writestr(name, content, compress_type=zipfile.ZIP_DEFLATED)
        return epub_path

    def test_epub_splits_into_chapters(self, sample_epub):
        """EPUB 应按 spine 拆为多个 PageContent（每章一页）。"""
        from parsers.epub_parser import EPUBParser

        parser = EPUBParser()
        pages = parser.parse(sample_epub)
        assert len(pages) == 2, f"期望 2 章，实际 {len(pages)}"

    def test_epub_chapter_title_resolved(self, sample_epub):
        """element metadata.chapter_title 应填上 NCX navLabel。"""
        from parsers.epub_parser import EPUBParser

        parser = EPUBParser()
        pages = parser.parse(sample_epub)
        titles = []
        for pg in pages:
            for elem in pg.elements:
                t = (elem.metadata or {}).get("chapter_title")
                if t is not None:
                    titles.append(t)
        assert "第一章 引言" in titles
        assert "第二章 方法" in titles


class TestR11PptxAnimations:
    """v0.11.0/B6: PPTX 动画顺序 + 讲者备注。"""

    # v0.13.0: python-pptx 在某些 venv 缺失 → importorskip 防 ERROR
    pptx = pytest.importorskip("pptx")

    @pytest.fixture
    def sample_pptx(self, tmp_path):
        """生成带 1 张幻灯片 + 1 备注 + 1 动画的 PPTX。"""
        from pptx import Presentation
        from pptx.util import Inches
        from lxml import etree

        pptx_path = tmp_path / "mini.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # title slide layout
        # title
        slide.shapes.title.text = "测试标题"
        # body
        from pptx.util import Pt as _Pt

        body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body is not None:
            body.text = "测试正文"
        # 备注
        notes_slide = slide.notes_slide
        if notes_slide:
            notes_slide.notes_text_frame.text = "这是讲者备注"
        # 手动注入一个 p:timing 节点（python-pptx 不直接暴露添加动画 API，用 XML 注入）
        slide_elem = slide._element
        # 找 sld 节点并在其末尾插入 cSld/timing
        timing_xml = (
            '<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:tnLst>'
            '<p:par>'
            '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"/>'
            "</p:par>"
            '<p:par>'
            '<p:cTn id="2" dur="2000" fill="hold" nodeType="clickEffect">'
            '<p:stCell val="1"/>'
            '<p:spTgt>'
            '<p:tgtEl spid="' + str(slide.shapes.title._element.get("id", "1")) + '"/>'
            "</p:spTgt>"
            "</p:cTn>"
            "<p:animMotion origin=\"center\" path=\"M 0,0 L 0,0\" dur=\"500ms\"/>"
            "</p:par>"
            "</p:tnLst>"
            "</p:timing>"
        )
        timing_elem = etree.fromstring(timing_xml)
        # timing 是 cSld 的兄弟节点（在 PresentationML schema 里）
        csld = slide_elem.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld"
        )
        if csld is not None:
            csld.addnext(timing_elem)

        prs.save(pptx_path)
        return pptx_path

    def test_pptx_speaker_notes_extracted(self, sample_pptx):
        """讲者备注应作为 elementType=note 元素抽出。"""
        from parsers.pptx_parser import PPTXParser

        pages = PPTXParser().parse(sample_pptx)
        notes = [e for pg in pages for e in pg.elements if e.elementType == "note"]
        assert len(notes) >= 1
        assert "讲者备注" in notes[0].content

    def test_pptx_animation_order_extracted(self, sample_pptx):
        """p:timing 里的动画应抽出为 metadata.animations 列表。"""
        from parsers.pptx_parser import PPTXParser

        pages = PPTXParser().parse(sample_pptx)
        # 第一页应至少有一条动画
        meta = pages[0].metadata or {}
        animations = meta.get("animations", [])
        assert meta.get("animations_count") == len(animations)
        # 至少有一个 animMotion
        assert any(a.get("effect_type") == "animMotion" for a in animations)
        # index 递增
        indices = [a["index"] for a in animations]
        assert indices == sorted(indices) and len(set(indices)) == len(indices)


class TestR12Diff:
    """v0.12.0/B10: diff 子命令 LCS 对比。"""

    def test_diff_simple_versions(self, tmp_path):
        """两份版本 → additions/deletions/unchanged + diff 预览。"""
        a = tmp_path / "a.txt"
        b = tmp_path / "b.txt"
        a.write_text("line1\nline2\nline3\n", encoding="utf-8")
        b.write_text("line1\nline2-changed\nline3\nline4\n", encoding="utf-8")
        payload, code = run_cli("diff", str(a), str(b))
        assert payload["ok"] is True
        d = payload["data"]
        assert d["additions"] >= 1
        assert d["deletions"] >= 1
        assert d["unchanged_count"] >= 2
        # 相似度 0~1
        assert 0 <= d["similarity"] <= 1
        # diff 预览包含 +/- 行
        assert "+" in d["diff_preview"]
        assert "-" in d["diff_preview"]
        assert code == 0

    def test_diff_identical_files(self, tmp_path):
        """相同文件 → 0 增 0 删，全部未变。"""
        a = tmp_path / "a.txt"
        b = tmp_path / "b.txt"
        a.write_text("same\ncontent\n", encoding="utf-8")
        b.write_text("same\ncontent\n", encoding="utf-8")
        payload, code = run_cli("diff", str(a), str(b))
        assert payload["ok"] is True
        d = payload["data"]
        assert d["additions"] == 0
        assert d["deletions"] == 0
        assert d["unchanged_count"] >= 2
        assert code == 0

    def test_diff_missing_file(self, tmp_path):
        """源不存在 → file_not_found。"""
        a = tmp_path / "a.txt"
        a.write_text("x\n", encoding="utf-8")
        payload, code = run_cli("diff", str(a), str(tmp_path / "missing.txt"))
        assert payload["ok"] is False
        assert payload["error"]["kind"] == "file_not_found"
        assert code != 0

    def test_diff_pdf_support(self, tmp_path):
        """PDF 文件也可 diff（走 translate 中间转换）。"""
        target = FIXTURES / "complex_test.pdf"
        if not target.exists():
            pytest.skip("fixture 缺失")
        payload, code = run_cli("diff", str(target), str(target), "--format", "text")
        assert payload["ok"] is True
        # 相同文件 → 0 增 0 删
        assert payload["data"]["additions"] == 0
        assert payload["data"]["deletions"] == 0
        assert code == 0
