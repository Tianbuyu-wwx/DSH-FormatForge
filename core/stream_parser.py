"""
流式解析器模块
支持大文件分块读取和内存高效处理
"""
import io
import logging
from abc import ABC, abstractmethod
from typing import Iterator, List, Optional, Dict, Any, BinaryIO
from pathlib import Path
from dataclasses import dataclass

logger = logging.getLogger("stream_parser")


@dataclass
class TextChunk:
    """文本分块"""
    chunk_id: int
    content: str
    start_pos: int
    end_pos: int
    is_complete_sentence: bool = False
    metadata: Optional[Dict[str, Any]] = None


class StreamParser(ABC):
    """流式解析器抽象基类"""

    def __init__(self, chunk_size: int = 8192, overlap: int = 200):
        self.chunk_size = chunk_size
        self.overlap = overlap

    @abstractmethod
    def can_handle(self, file_path: Path) -> bool:
        """是否能处理该文件"""
        pass

    @abstractmethod
    def parse_stream(self, file_path: Path) -> Iterator[TextChunk]:
        """流式解析文件，返回文本分块"""
        pass


class TextStreamParser(StreamParser):
    """文本文件流式解析器"""

    def can_handle(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in {".txt", ".md", ".csv", ".log", ".json", ".xml", ".yaml", ".yml"}

    def parse_stream(self, file_path: Path) -> Iterator[TextChunk]:
        """按语义边界分块读取文本"""
        chunk_id = 0
        buffer = ""
        start_pos = 0

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            while True:
                data = f.read(self.chunk_size)
                if not data:
                    break

                buffer += data

                # 找到最后一个句子边界
                boundary = self._find_sentence_boundary(buffer)
                if boundary > 0:
                    content = buffer[:boundary]
                    yield TextChunk(
                        chunk_id=chunk_id,
                        content=content,
                        start_pos=start_pos,
                        end_pos=start_pos + len(content),
                        is_complete_sentence=True
                    )
                    chunk_id += 1
                    start_pos += len(content)
                    # 保留重叠部分
                    buffer = buffer[boundary - self.overlap:boundary] if self.overlap < boundary else buffer[boundary:]

        # 处理剩余内容
        if buffer.strip():
            yield TextChunk(
                chunk_id=chunk_id,
                content=buffer.strip(),
                start_pos=start_pos,
                end_pos=start_pos + len(buffer.strip()),
                is_complete_sentence=False
            )

    def _find_sentence_boundary(self, text: str) -> int:
        """查找句子边界（句号、换行等）"""
        boundaries = []
        for marker in ["\n\n", ". ", "。", "! ", "? ", "\n"]:
            pos = text.rfind(marker, self.chunk_size // 2)
            if pos > 0:
                boundaries.append(pos + len(marker))

        return max(boundaries) if boundaries else 0


class BinaryStreamParser(StreamParser):
    """二进制文件流式解析器（用于PDF、DOCX等）"""

    def __init__(self, chunk_size: int = 65536, overlap: int = 0):
        super().__init__(chunk_size, overlap)
        self._parsers = {
            ".pdf": self._parse_pdf_stream,
            ".docx": self._parse_docx_stream,
            ".pptx": self._parse_pptx_stream,
            ".xlsx": self._parse_xlsx_stream,
        }

    def can_handle(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in self._parsers

    def parse_stream(self, file_path: Path) -> Iterator[TextChunk]:
        """根据文件类型选择对应的流式解析"""
        suffix = file_path.suffix.lower()
        parser = self._parsers.get(suffix)
        if parser:
            yield from parser(file_path)
        else:
            logger.warning(f"不支持的文件类型: {suffix}")

    def _parse_pdf_stream(self, file_path: Path) -> Iterator[TextChunk]:
        """流式解析PDF（逐页处理）"""
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(file_path))
            total_pages = len(reader.pages)

            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    yield TextChunk(
                        chunk_id=i,
                        content=text,
                        start_pos=i,
                        end_pos=i + 1,
                        metadata={"page": i + 1, "total_pages": total_pages}
                    )
        except Exception as e:
            logger.error(f"PDF流式解析失败: {e}")
            raise

    def _parse_docx_stream(self, file_path: Path) -> Iterator[TextChunk]:
        """流式解析DOCX（按段落分块）"""
        try:
            from docx import Document
            doc = Document(str(file_path))
            chunk_id = 0
            buffer = []
            current_size = 0

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                buffer.append(text)
                current_size += len(text)

                if current_size >= self.chunk_size:
                    content = "\n".join(buffer)
                    yield TextChunk(
                        chunk_id=chunk_id,
                        content=content,
                        start_pos=chunk_id,
                        end_pos=chunk_id + 1
                    )
                    chunk_id += 1
                    buffer = []
                    current_size = 0

            if buffer:
                content = "\n".join(buffer)
                yield TextChunk(
                    chunk_id=chunk_id,
                    content=content,
                    start_pos=chunk_id,
                    end_pos=chunk_id + 1
                )
        except Exception as e:
            logger.error(f"DOCX流式解析失败: {e}")
            raise

    def _parse_pptx_stream(self, file_path: Path) -> Iterator[TextChunk]:
        """流式解析PPTX（逐幻灯片处理）"""
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            total_slides = len(prs.slides)

            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)

                content = "\n".join(texts)
                if content.strip():
                    yield TextChunk(
                        chunk_id=i,
                        content=content,
                        start_pos=i,
                        end_pos=i + 1,
                        metadata={"slide": i + 1, "total_slides": total_slides}
                    )
        except Exception as e:
            logger.error(f"PPTX流式解析失败: {e}")
            raise

    def _parse_xlsx_stream(self, file_path: Path) -> Iterator[TextChunk]:
        """流式解析XLSX（按工作表分块）"""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), read_only=True, data_only=True)

            chunk_id = 0
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                row_count = 0

                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    rows.append(row_text)
                    row_count += 1

                    if row_count >= 1000:  # 每1000行一个分块
                        content = f"## 工作表: {sheet_name}\n" + "\n".join(rows)
                        yield TextChunk(
                            chunk_id=chunk_id,
                            content=content,
                            start_pos=chunk_id,
                            end_pos=chunk_id + 1,
                            metadata={"sheet": sheet_name, "rows": row_count}
                        )
                        chunk_id += 1
                        rows = []
                        row_count = 0

                if rows:
                    content = f"## 工作表: {sheet_name}\n" + "\n".join(rows)
                    yield TextChunk(
                        chunk_id=chunk_id,
                        content=content,
                        start_pos=chunk_id,
                        end_pos=chunk_id + 1,
                        metadata={"sheet": sheet_name, "rows": row_count}
                    )
                    chunk_id += 1
        except Exception as e:
            logger.error(f"XLSX流式解析失败: {e}")
            raise


class StreamParserManager:
    """流式解析器管理器"""

    STREAMING_THRESHOLD = 10 * 1024 * 1024  # 10MB 阈值

    def __init__(self):
        self._parsers: List[StreamParser] = [
            TextStreamParser(),
            BinaryStreamParser(),
        ]

    def should_use_streaming(self, file_path: Path) -> bool:
        """判断是否应该使用流式处理"""
        if not file_path.exists():
            return False
        return file_path.stat().st_size > self.STREAMING_THRESHOLD

    def parse(self, file_path: Path) -> Iterator[TextChunk]:
        """自动选择合适的解析器进行流式处理"""
        path = Path(file_path)

        for parser in self._parsers:
            if parser.can_handle(path):
                logger.info(f"使用 {parser.__class__.__name__} 流式解析: {path}")
                yield from parser.parse_stream(path)
                return

        raise ValueError(f"无法流式解析文件: {path}")

    def parse_with_progress(self, file_path: Path) -> Iterator[TextChunk]:
        """带进度信息的流式解析"""
        total_size = Path(file_path).stat().st_size
        processed_size = 0

        for chunk in self.parse(file_path):
            processed_size += len(chunk.content.encode("utf-8"))
            chunk.metadata = chunk.metadata or {}
            chunk.metadata["progress"] = min(processed_size / total_size, 1.0)
            chunk.metadata["total_size"] = total_size
            yield chunk


# 全局管理器
stream_parser_manager = StreamParserManager()
